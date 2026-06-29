import hashlib
import io
import json
import csv
import os
import re
from datetime import datetime
from html import unescape
from html.parser import HTMLParser
from typing import Any, Callable, Dict, List, Optional
from uuid import uuid4

import yaml

from .config import settings
from .graph import create_document_graph, get_neo4j_driver
from .text_utils import chunk_text, normalize_text


def _sha1(text: str) -> str:
    return hashlib.sha1(text.encode("utf-8")).hexdigest()


class _HTMLTextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts: List[str] = []

    def handle_data(self, data: str) -> None:
        self.parts.append(data)

    def handle_entityref(self, name: str) -> None:
        self.parts.append(unescape(f"&{name};"))

    def handle_charref(self, name: str) -> None:
        self.parts.append(unescape(f"&#{name};"))

    def get_text(self) -> str:
        return " ".join(self.parts)


def extract_text_from_html(raw: str) -> str:
    parser = _HTMLTextExtractor()
    parser.feed(raw)
    parser.close()
    return normalize_text(parser.get_text())


def extract_text_from_json(raw: str) -> str:
    try:
        parsed = json.loads(raw)
        return normalize_text(json.dumps(parsed, indent=2, ensure_ascii=False))
    except Exception:
        return normalize_text(raw)


def extract_text_from_yaml(raw: str) -> str:
    try:
        parsed = yaml.safe_load(raw)
        return normalize_text(yaml.dump(parsed, sort_keys=False, allow_unicode=True))
    except Exception:
        return normalize_text(raw)


def extract_text_from_csv(raw: str) -> str:
    try:
        rows = []
        reader = csv.reader(io.StringIO(raw))
        for row in reader:
            rows.append(" | ".join(row))
        return normalize_text("\n".join(rows))
    except Exception:
        return normalize_text(raw)


def extract_text_from_xml(raw: str) -> str:
    try:
        import xml.etree.ElementTree as ET
        root = ET.fromstring(raw)

        parts: List[str] = []

        def walk(node: Any) -> None:
            if node.text and node.text.strip():
                parts.append(node.text.strip())
            for child in node:
                walk(child)
            if node.tail and node.tail.strip():
                parts.append(node.tail.strip())

        walk(root)
        return normalize_text("\n".join(parts))
    except Exception:
        return normalize_text(raw)


def extract_text_from_docx(docx_bytes: bytes) -> str:
    try:
        from docx import Document
        document = Document(io.BytesIO(docx_bytes))
        paragraphs = [p.text for p in document.paragraphs if p.text]
        table_rows: List[str] = []
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text for cell in row.cells if cell.text]
                if cells:
                    table_rows.append(" | ".join(cells))
        text = "\n".join(paragraphs + table_rows)
        return normalize_text(text)
    except ImportError:
        return _extract_docx_via_zip(docx_bytes)
    except Exception:
        fallback = _extract_docx_via_zip(docx_bytes)
        if fallback:
            return fallback
        raise


def _extract_docx_via_zip(docx_bytes: bytes) -> str:
    import zipfile
    with zipfile.ZipFile(io.BytesIO(docx_bytes)) as archive:
        xml = archive.read("word/document.xml").decode("utf-8", errors="ignore")
    # Prefer Word XML text nodes when present.
    text_nodes = re.findall(r"<w:t[^>]*>([^<]*)</w:t>", xml)
    if text_nodes:
        cleaned = normalize_text("\n".join(text_nodes))
        if cleaned:
            return cleaned
    text = re.sub(r"<w:tab[^/]*/>", "\t", xml)
    text = re.sub(r"</w:p>", "\n", text)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+\n", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    cleaned = normalize_text(text)
    if not cleaned:
        raise RuntimeError("Could not read text from this Word document.")
    return cleaned


def extract_text_from_doc_bytes(file_bytes: bytes, filename: str = "") -> str:
    """Extract text from .docx or legacy .doc when possible."""
    lower = (filename or "").lower()
    if file_bytes[:2] == b"PK":
        return extract_text_from_docx(file_bytes)
    if lower.endswith(".docx"):
        return extract_text_from_docx(file_bytes)
    if lower.endswith(".doc"):
        # Legacy binary .doc is not reliably supported; try OOXML fallback anyway.
        try:
            return extract_text_from_docx(file_bytes)
        except Exception as exc:
            raise RuntimeError(
                "This is an older .doc file. Please save it as .docx or .pdf and upload again."
            ) from exc
    return extract_text_from_docx(file_bytes)


def extract_text_from_pptx(pptx_bytes: bytes) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise RuntimeError("python-pptx is required to ingest .pptx files. Install it via requirements.") from exc

    presentation = Presentation(io.BytesIO(pptx_bytes))
    slides: List[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slides.append(shape.text)
    return normalize_text("\n".join(slides))


def extract_text_from_file(filename: str, file_bytes: bytes) -> str:
    ext = os.path.splitext(filename.lower())[1]
    raw_text = file_bytes.decode("utf-8", errors="ignore")

    if ext in {".txt", ".md", ".log"}:
        return normalize_text(raw_text)
    if ext in {".json"}:
        return extract_text_from_json(raw_text)
    if ext in {".yaml", ".yml"}:
        return extract_text_from_yaml(raw_text)
    if ext in {".csv"}:
        return extract_text_from_csv(raw_text)
    if ext in {".xml"}:
        return extract_text_from_xml(raw_text)
    if ext in {".html", ".htm"}:
        return extract_text_from_html(raw_text)
    if ext == ".docx" or ext == ".doc":
        return extract_text_from_doc_bytes(file_bytes, filename)
    if ext == ".pptx":
        return extract_text_from_pptx(file_bytes)

    raise RuntimeError(f"Unsupported text document type: {ext}")


async def ingest_text_file_to_mongo(
    db,
    filename: str,
    file_bytes: bytes,
    embed_fn: Callable[[str], Optional[List[float]]],
    metadata: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    metadata = metadata or {}
    doc_id = str(uuid4())
    text = extract_text_from_file(filename, file_bytes)
    if not text:
        raise RuntimeError(f"Text extraction failed: no content found in {filename}")

    chunks = chunk_text(text, settings.chunk_size, settings.chunk_overlap)
    if not chunks:
        raise RuntimeError(f"Text file has no extractable content: {filename}")

    chunk_docs: List[Dict[str, Any]] = []
    for idx, chunk in enumerate(chunks):
        emb = embed_fn(chunk)
        chunk_doc = {
            "doc_id": doc_id,
            "source": filename,
            "source_type": "text",
            "page": 1,
            "chunk_index": idx,
            "chunk_hash": _sha1(chunk),
            "text": chunk,
            "metadata": metadata,
            "created_at": datetime.utcnow(),
        }
        if emb is not None:
            chunk_doc["embedding"] = emb
        else:
            chunk_doc["embedding"] = None
            print(f"WARNING: Text chunk {idx} in {filename} has no embedding")
        chunk_docs.append(chunk_doc)

    await db[settings.chunks_collection].insert_many(chunk_docs)
    await db[settings.docs_collection].insert_one(
        {
            "doc_id": doc_id,
            "filename": filename,
            "pages": 1,
            "chunks": len(chunk_docs),
            "metadata": metadata,
            "created_at": datetime.utcnow(),
        }
    )

    driver = get_neo4j_driver()
    if driver is not None:
        try:
            await create_document_graph(driver, doc_id, filename, chunk_docs, metadata)
        except Exception as graph_err:
            print("NEO4J TEXT INGEST ERROR:", graph_err)

    return {
        "doc_id": doc_id,
        "filename": filename,
        "pages_indexed": 1,
        "chunks_indexed": len(chunk_docs),
    }


async def ingest_docx_to_mongo(
    db,
    filename: str,
    docx_bytes: bytes,
    embed_fn: Callable[[str], Optional[List[float]]],
    metadata: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    metadata = metadata or {}
    text = extract_text_from_docx(docx_bytes)
    return await ingest_text_file_to_mongo(db, filename, text.encode("utf-8"), embed_fn, metadata)


async def ingest_pptx_to_mongo(
    db,
    filename: str,
    pptx_bytes: bytes,
    embed_fn: Callable[[str], Optional[List[float]]],
    metadata: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    metadata = metadata or {}
    text = extract_text_from_pptx(pptx_bytes)
    return await ingest_text_file_to_mongo(db, filename, text.encode("utf-8"), embed_fn, metadata)


def extract_pdf_pages(pdf_bytes: bytes) -> List[Dict[str, Any]]:
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception as e:
        raise RuntimeError("pypdf is not installed. Please install requirements.") from e
    reader = PdfReader(io.BytesIO(pdf_bytes))
    pages: List[Dict[str, Any]] = []
    total_raw_chars = 0
    
    print(f"📄 PDF: Extracting from {len(reader.pages)} pages...")
    for i, page in enumerate(reader.pages):
        raw = page.extract_text() or ""
        total_raw_chars += len(raw)
        txt = normalize_text(raw)
        if txt:
            pages.append({"page": i + 1, "text": txt})
            print(f"  ✓ Page {i + 1}: {len(raw)} raw chars → {len(txt)} normalized chars")
        else:
            print(f"  ⚠️ Page {i + 1}: No text extracted ({len(raw)} raw chars)")
    
    print(f"✓ PDF extraction: {len(pages)}/{len(reader.pages)} pages with text, {total_raw_chars} total raw chars")
    return pages


async def ingest_pdf_to_mongo(
    db,
    filename: str,
    pdf_bytes: bytes,
    embed_fn: Callable[[str], Optional[List[float]]],
    metadata: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    metadata = metadata or {}
    doc_id = str(uuid4())
    pages = extract_pdf_pages(pdf_bytes)

    if not pages:
        raise RuntimeError(f"PDF parsing failed: no text extracted from {filename}")

    chunk_docs = []
    total_chunks = 0
    chunks_without_embedding = 0

    print(f"📋 INGEST: PDF has {len(pages)} pages, metadata={metadata}")

    for page_obj in pages:
        page_no = page_obj["page"]
        text = page_obj["text"]
        chunks = chunk_text(text, settings.chunk_size, settings.chunk_overlap)

        for idx, chunk in enumerate(chunks):
            emb = embed_fn(chunk)
            chunk_doc = {
                "doc_id": doc_id,
                "source": filename,
                "source_type": "pdf",
                "page": page_no,
                "chunk_index": idx,
                "chunk_hash": _sha1(chunk),
                "text": chunk,
                "metadata": metadata,
                "created_at": datetime.utcnow(),
            }
            
            if emb is not None:
                chunk_doc["embedding"] = emb
                chunk_docs.append(chunk_doc)
                total_chunks += 1
            else:
                # Store chunk even without embedding for fallback retrieval
                chunk_doc["embedding"] = None
                chunk_docs.append(chunk_doc)
                chunks_without_embedding += 1
                total_chunks += 1
                print(f"WARNING: Chunk {idx} on page {page_no} has no embedding")

    if not chunk_docs:
        raise RuntimeError(f"PDF has no extractable content: {filename}")

    print(f"💾 INGEST: Storing {len(chunk_docs)} chunks to collection '{settings.chunks_collection}'")
    await db[settings.chunks_collection].insert_many(chunk_docs)
    print(f"✓ Ingested {filename}: {total_chunks} chunks ({chunks_without_embedding} without embeddings)")

    await db[settings.docs_collection].insert_one(
        {
            "doc_id": doc_id,
            "filename": filename,
            "pages": len(pages),
            "chunks": total_chunks,
            "metadata": metadata,
            "created_at": datetime.utcnow(),
        }
    )
    print(f"✓ Stored document record to collection '{settings.docs_collection}'")

    driver = get_neo4j_driver()
    if driver is not None:
        try:
            await create_document_graph(driver, doc_id, filename, chunk_docs, metadata)
        except Exception as graph_err:
            print("NEO4J INGEST ERROR:", graph_err)

    return {
        "doc_id": doc_id,
        "filename": filename,
        "pages_indexed": len(pages),
        "chunks_indexed": total_chunks,
    }


def _extract_image_text(image_bytes: bytes) -> str:
    # Best-effort OCR. Works when pillow+pytesseract are available and tesseract is installed.
    try:
        import pytesseract  # type: ignore
        from PIL import Image  # type: ignore
    except Exception:
        return ""

    try:
        img = Image.open(io.BytesIO(image_bytes))
        text = pytesseract.image_to_string(img) or ""
        return normalize_text(text)
    except Exception:
        return ""


async def ingest_image_to_mongo(
    db,
    filename: str,
    image_bytes: bytes,
    embed_fn: Callable[[str], Optional[List[float]]],
    metadata: Optional[Dict[str, Any]] = None,
    vision_fn: Optional[Callable[[bytes, str], str]] = None,
) -> Dict[str, Any]:
    metadata = metadata or {}
    doc_id = str(uuid4())

    text = ""
    if vision_fn is not None:
        try:
            text = (vision_fn(image_bytes, filename) or "").strip()
            if text:
                print(f"✓ IMAGE VISION: Extracted description for {filename} ({len(text)} chars)")
        except Exception as exc:
            print(f"IMAGE VISION FN ERROR ({filename}): {exc}")

    if not text:
        text = _extract_image_text(image_bytes)

    if not text:
        text = f"Image uploaded: {filename}. No readable text could be extracted."

    emb = embed_fn(text)
    total_chunks = 0
    chunk_docs = []
    
    chunk_doc = {
        "doc_id": doc_id,
        "source": filename,
        "source_type": "image",
        "page": 1,
        "chunk_index": 0,
        "chunk_hash": _sha1(text),
        "text": text,
        "metadata": metadata,
        "created_at": datetime.utcnow(),
    }
    
    if emb is not None:
        chunk_doc["embedding"] = emb
    else:
        chunk_doc["embedding"] = None
        print(f"WARNING: Image {filename} has no embedding")
    
    chunk_docs.append(chunk_doc)
    await db[settings.chunks_collection].insert_one(chunk_docs[0])
    total_chunks = 1
    print(f"✓ Ingested {filename}: 1 chunk ({0 if emb else 1} without embedding)")

    await db[settings.docs_collection].insert_one(
        {
            "doc_id": doc_id,
            "filename": filename,
            "pages": 1,
            "chunks": total_chunks,
            "metadata": metadata,
            "created_at": datetime.utcnow(),
        }
    )

    driver = get_neo4j_driver()
    if driver is not None and chunk_docs:
        try:
            await create_document_graph(driver, doc_id, filename, chunk_docs, metadata)
        except Exception as graph_err:
            print("NEO4J IMAGE INGEST ERROR:", graph_err)

    return {
        "doc_id": doc_id,
        "filename": filename,
        "pages_indexed": 1,
        "chunks_indexed": total_chunks,
    }

